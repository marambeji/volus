from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define Novelus brand colors (purple theme from logo)
PURPLE_PRIMARY = RGBColor(134, 59, 255)   # #863bff from logo
PURPLE_DARK = RGBColor(126, 20, 255)      # #7e14ff from logo
PURPLE_LIGHT = RGBColor(237, 230, 255)    # #ede6ff from logo
BLUE_ACCENT = RGBColor(71, 191, 255)      # #47bfff from logo
GRAY_DARK = RGBColor(51, 51, 51)
GRAY_LIGHT = RGBColor(128, 128, 128)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 139, 34)
RED = RGBColor(220, 53, 69)

def add_logo(slide, size=0.6):
    """Add Novelus logo to slide"""
    # Try the actual Novelus logo first
    logo_path = "C:/Users/user/Downloads/NOV-Logo-IN.svg"
    if not os.path.exists(logo_path):
        logo_path = "frontend/public/favicon.svg"

    if os.path.exists(logo_path):
        try:
            logo_left = Inches(0.4)
            logo_top = Inches(0.3)
            logo_pic = slide.shapes.add_picture(logo_path, logo_left, logo_top,
                                                height=Inches(size))
        except:
            # If SVG doesn't work, create a placeholder
            add_logo_placeholder(slide, size)
    else:
        add_logo_placeholder(slide, size)

def add_logo_placeholder(slide, size=0.6):
    """Add logo placeholder if SVG cannot be loaded"""
    logo = slide.shapes.add_shape(1, Inches(0.4), Inches(0.3),
                                  Inches(size), Inches(size))
    logo.fill.solid()
    logo.fill.fore_color.rgb = PURPLE_PRIMARY
    logo.line.fill.background()

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple accent bar
    accent = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.25))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE_PRIMARY
    accent.line.fill.background()

    # Add logo
    try:
        logo_path = "C:/Users/user/Downloads/NOV-Logo-IN.svg"
        if not os.path.exists(logo_path):
            logo_path = "frontend/public/favicon.svg"
        if os.path.exists(logo_path):
            logo_left = (prs.slide_width - Inches(1.5)) / 2
            slide.shapes.add_picture(logo_path, logo_left, Inches(1.5), height=Inches(1.5))
    except:
        pass

    # Novelus text
    logo_text = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(0.6))
    text_frame = logo_text.text_frame
    text_frame.text = "NOVELUS"
    text_frame.paragraphs[0].font.size = Pt(48)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = PURPLE_PRIMARY
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = GRAY_DARK
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.word_wrap = True

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = GRAY_LIGHT
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "August 2026"
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].font.color.rgb = GRAY_LIGHT
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_toc_slide(prs):
    """Add table of contents slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple header
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE_PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Table of Contents"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # TOC items
    toc_items = [
        "1. Introduction",
        "2. Problem Statement",
        "3. Existing System Analysis",
        "4. Proposed Solution",
        "5. Functional Requirements",
        "6. Technical Requirements",
        "7. Added Value and Innovation",
        "8. Accessibility",
        "9. Performance",
        "10. Demo and Video",
        "11. GitHub Repository and LinkedIn",
        "12. Conclusion"
    ]

    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(7), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(toc_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY_DARK
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_section_divider(prs, section_number, section_title):
    """Add section divider"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Purple background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_PRIMARY
    background.line.fill.background()

    # Section number
    number_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.8))
    number_frame = number_box.text_frame
    number_frame.text = f"{section_number:02d}"
    number_frame.paragraphs[0].font.size = Pt(72)
    number_frame.paragraphs[0].font.bold = True
    number_frame.paragraphs[0].font.color.rgb = PURPLE_LIGHT
    number_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.paragraphs[0].font.size = Pt(42)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, items, item_color=GRAY_DARK):
    """Add content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple header
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE_PRIMARY
    header.line.fill.background()

    # Add small logo
    try:
        logo_path = "C:/Users/user/Downloads/NOV-Logo-IN.svg"
        if not os.path.exists(logo_path):
            logo_path = "frontend/public/favicon.svg"
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(9.3), Inches(0.25), height=Inches(0.5))
    except:
        pass

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = item_color
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(6)

def add_comparison_slide(prs, title, old_items, new_items):
    """Add before/after comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple header
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE_PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.4))
    left_header_frame = left_header.text_frame
    left_header_frame.text = "Old PowerApps Portal"
    left_header_frame.paragraphs[0].font.size = Pt(20)
    left_header_frame.paragraphs[0].font.bold = True
    left_header_frame.paragraphs[0].font.color.rgb = RED

    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.3), Inches(5.3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(old_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = f"✗ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY_DARK
        p.space_before = Pt(8)

    # Divider
    divider = slide.shapes.add_shape(1, Inches(4.9), Inches(1.3), Inches(0.05), Inches(5.9))
    divider.fill.solid()
    divider.fill.fore_color.rgb = PURPLE_LIGHT
    divider.line.fill.background()

    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.4))
    right_header_frame = right_header.text_frame
    right_header_frame.text = "New Novelus HR Portal"
    right_header_frame.paragraphs[0].font.size = Pt(20)
    right_header_frame.paragraphs[0].font.bold = True
    right_header_frame.paragraphs[0].font.color.rgb = GREEN

    # Right content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.3), Inches(5.3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, item in enumerate(new_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = f"✓ {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY_DARK
        p.space_before = Pt(8)

def add_screenshot_slide(prs, title, description, image_path=None):
    """Add slide with screenshot"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple header
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PURPLE_PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE

    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.4))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.paragraphs[0].font.size = Pt(16)
    desc_frame.paragraphs[0].font.color.rgb = GRAY_DARK

    # Add image if exists
    if image_path and os.path.exists(image_path):
        try:
            # Calculate dimensions to fit in available space
            img_left = Inches(0.8)
            img_top = Inches(1.8)
            img_width = Inches(8.4)
            img_height = Inches(5)

            slide.shapes.add_picture(image_path, img_left, img_top,
                                    width=img_width, height=img_height)
        except Exception as e:
            # If image fails, show placeholder
            placeholder = slide.shapes.add_shape(1, Inches(0.8), Inches(1.8),
                                               Inches(8.4), Inches(5))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = PURPLE_LIGHT
            placeholder.line.color.rgb = PURPLE_PRIMARY
            placeholder.line.width = Pt(2)

            ph_text = placeholder.text_frame
            ph_text.text = f"[Screenshot: {title}]"
            ph_text.paragraphs[0].font.size = Pt(20)
            ph_text.paragraphs[0].font.color.rgb = PURPLE_PRIMARY
            ph_text.paragraphs[0].alignment = PP_ALIGN.CENTER
            ph_text.vertical_anchor = 1
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(1, Inches(0.8), Inches(1.8),
                                           Inches(8.4), Inches(5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = PURPLE_LIGHT
        placeholder.line.color.rgb = PURPLE_PRIMARY
        placeholder.line.width = Pt(2)

        ph_text = placeholder.text_frame
        ph_text.text = f"[Screenshot: {title}]"
        ph_text.paragraphs[0].font.size = Pt(20)
        ph_text.paragraphs[0].font.color.rgb = PURPLE_PRIMARY
        ph_text.paragraphs[0].alignment = PP_ALIGN.CENTER
        ph_text.vertical_anchor = 1

def add_closing_slide(prs):
    """Add closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # Purple accent
    accent = slide.shapes.add_shape(1, 0, Inches(2.8), prs.slide_width, Inches(2.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE_LIGHT
    accent.line.fill.background()

    # Logo
    try:
        logo_path = "C:/Users/user/Downloads/NOV-Logo-IN.svg"
        if not os.path.exists(logo_path):
            logo_path = "frontend/public/favicon.svg"
        if os.path.exists(logo_path):
            logo_left = (prs.slide_width - Inches(1.2)) / 2
            slide.shapes.add_picture(logo_path, logo_left, Inches(1.5), height=Inches(1.2))
    except:
        pass

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PURPLE_PRIMARY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "HR Portal Modernization Project"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = GRAY_DARK
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============= CREATE PRESENTATION =============

print("Creating presentation...")

# Title Slide
add_title_slide(prs, "HR Portal Modernization", "Leave Management System")

# Table of Contents
add_toc_slide(prs)

# 1. Introduction
add_section_divider(prs, 1, "Introduction")
add_content_slide(prs, "Project Overview", [
    "Project: HR Leave Portal Modernization for Novelus",
    "Goal: Replace legacy PowerApps system with modern web application",
    "Timeline: Development completed August 2026",
    "Technology: React 19, TypeScript, Vite, NestJS, PostgreSQL",
    "Users: Employees, Managers, and HR Admins"
])

# 2. Problem Statement
add_section_divider(prs, 2, "Problem Statement")
add_content_slide(prs, "Challenges with Current System", [
    "Old PowerApps interface is difficult to use",
    "Poor mobile device support",
    "Manual balance tracking causes errors",
    "No real-time updates",
    "Limited reporting capabilities",
    "Difficult to maintain"
], RED)

# 3. Existing System Analysis
add_section_divider(prs, 3, "Existing System Analysis")
add_content_slide(prs, "Old PowerApps Portal", [
    "URL: https://apps.powerapps.com/play/...",
    "Platform: Microsoft PowerApps",
    "Issues: Slow performance, complex navigation",
    "Limitations: Basic features, no mobile optimization",
    "Maintenance: Difficult to update and extend"
])

# 4. Proposed Solution
add_section_divider(prs, 4, "Proposed Solution")
add_content_slide(prs, "New Novelus HR Portal", [
    "URL: https://hr.novelus.com/dev/lr/login",
    "Modern web application built from ground up",
    "Mobile-responsive design",
    "Automated workflows and balance tracking",
    "Real-time updates across system",
    "Advanced reporting and analytics"
], GREEN)

add_comparison_slide(prs, "Before & After",
    ["Slow, outdated interface", "Poor mobile support", "Manual balance tracking",
     "No real-time updates", "Basic reporting", "Hard to maintain"],
    ["Modern, fast interface", "Mobile-first design", "Automated balance tracking",
     "Live updates", "Advanced analytics", "Easy to maintain"]
)

# 5. Functional Requirements
add_section_divider(prs, 5, "Functional Requirements")
add_content_slide(prs, "Employee Features", [
    "View leave balances",
    "Submit and track leave requests",
    "Cancel pending or approved requests",
    "View team calendar",
    "Receive notifications",
    "Access personal reports"
])

add_content_slide(prs, "Manager & HR Admin Features", [
    "Managers: Approve/reject requests, view team balances",
    "Managers: Access team calendar and reports",
    "HR Admins: Manage all employees and policies",
    "HR Admins: Configure approval workflows",
    "HR Admins: Adjust balances and run accruals",
    "HR Admins: Access company-wide reports and audit logs"
])

# 6. Technical Requirements
add_section_divider(prs, 6, "Technical Requirements")
add_content_slide(prs, "Technology Stack", [
    "Frontend: React 19 with TypeScript",
    "Build Tool: Vite",
    "Backend: NestJS (Node.js framework)",
    "Database: PostgreSQL with TypeORM",
    "API Documentation: Swagger",
    "Testing: Vitest, Jest, React Testing Library",
    "Authentication: Email-based login"
])

add_content_slide(prs, "System Architecture", [
    "Three-tier architecture",
    "Frontend: Single-page application (SPA)",
    "Backend: RESTful API with modular services",
    "Database: Relational model with migrations",
    "Scheduled jobs for automated tasks",
    "Role-based access control",
    "Complete audit logging"
])

# 7. Added Value and Innovation
add_section_divider(prs, 7, "Added Value and Innovation")
add_content_slide(prs, "Key Benefits", [
    "Reduced administrative burden on HR staff",
    "Automated workflows eliminate manual processes",
    "Improved accuracy in balance calculations",
    "Better employee and manager experience",
    "Data-driven insights from analytics",
    "Scalable system for company growth"
])

# 8. Accessibility
add_section_divider(prs, 8, "Accessibility")
add_content_slide(prs, "Device & Browser Compatibility", [
    "Responsive design for all screen sizes",
    "Mobile-optimized for phones and tablets",
    "Desktop-optimized for full functionality",
    "Compatible with Chrome, Firefox, Safari, Edge",
    "Works on Windows, macOS, Linux, iOS, Android",
    "No special software required"
])

# 9. Performance
add_section_divider(prs, 9, "Performance")
add_content_slide(prs, "Performance & Reliability", [
    "Optimized database queries",
    "Efficient API design",
    "Real-time updates without page refresh",
    "Comprehensive error handling",
    "Automated testing for reliability",
    "Built for scalability"
])

# 10. Demo and Video
add_section_divider(prs, 10, "Demo and Video")
add_screenshot_slide(prs, "Old PowerApps Portal",
                    "Legacy system with complex interface",
                    "frontend/public/login_bg.png")
add_screenshot_slide(prs, "New Novelus HR Portal - Login",
                    "Modern, clean login interface at https://hr.novelus.com/dev/lr/login",
                    "frontend/public/login_bg.png")
add_screenshot_slide(prs, "New Portal - Dashboard",
                    "Employee dashboard with balances and quick actions",
                    "frontend/src/assets/hero.png")

# 11. GitHub Repository and LinkedIn
add_section_divider(prs, 11, "GitHub Repository and LinkedIn")
add_content_slide(prs, "Project Resources", [
    "GitHub Repository:",
    "[Insert GitHub URL]",
    "",
    "LinkedIn:",
    "[Insert LinkedIn Profile]",
    "",
    "Documentation: Complete technical and business docs included",
    "Code: Full source code with comments and tests"
])

# 12. Conclusion
add_section_divider(prs, 12, "Conclusion")
add_content_slide(prs, "Project Summary", [
    "Successfully developed modern HR leave management system",
    "Replaces legacy PowerApps with user-friendly web application",
    "Provides comprehensive features for all user roles",
    "Built with modern, maintainable technology stack",
    "Includes complete documentation and testing",
    "Ready for deployment"
])

add_content_slide(prs, "Next Steps", [
    "1. User acceptance testing",
    "2. Data migration from old system",
    "3. User training sessions",
    "4. Phased rollout",
    "5. Production deployment",
    "6. Ongoing support and improvements"
])

# Closing
add_closing_slide(prs)

# Save
output_path = os.path.join(os.getcwd(), "Novelus_HR_Portal_Final.pptx")
prs.save(output_path)
print(f"Presentation created: {output_path}")
print(f"Total slides: {len(prs.slides)}")
